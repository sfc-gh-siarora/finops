import csv
import os
from datetime import date
from operator import concat

from pptx.util import Pt

import pandas as pd
from pptx.chart.data import CategoryChartData
from diskcache import Cache
from pptx import Presentation

from queries.billing.billing import BILLING
from snowhouse.browser_connection import set_connection
from snowhouse.dates import fetch_dates
from snowhouse.results import fetch_results_list_of_queries, fetch_results
from ue.uniteconomicswarehouselevel import get_deployment_zone_queries_wh_level, \
    get_deployment_zone_queries_quarter_level
from wh.warehouses_usage import get_warehouse_level_query

cache = Cache('./my_cache')

pd.set_option('display.max_rows', None)  # None means show all rows
pd.set_option('display.max_columns', None)
pd.set_option('display.float_format', '{:,.2f}'.format)


def cache_my_list(my_list, key):
    cache.set(key, my_list)


def retrieve_my_list(key):
    return cache.get(key, default=[])


def common_slides(company_name, formatted_date, prs, ve_name):
    replacemytext(prs, "<Company>", company_name)

    replacemytext(prs, "<Month>", formatted_date)
    replacemytext(prs, "<VE>", ve_name)


def format_number(number, decimal_places):
    if decimal_places == 0:
        # Round the number and convert to int to remove any decimal places
        return f"{int(round(number, decimal_places)):,}"
    else:
        # Format with specified number of decimal places and commas
        return f"{round(number, decimal_places):,.{decimal_places}f}"


def replacemytext(prs, search_str, repl_str):
    cnt = 0
    for slide in prs.slides:
        cnt += 1
        for shape in slide.shapes:
            if shape.has_text_frame:
                if (shape.text.find(search_str)) != -1:
                    text_frame = shape.text_frame
                    cur_text = text_frame.paragraphs[0].runs[0].text
                    new_text = cur_text.replace(str(search_str), str(repl_str))
                    text_frame.paragraphs[0].runs[0].text = new_text
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    for cell in row.cells:
                        if not cell.text_frame:
                            continue
                        for paragraph in cell.text_frame.paragraphs:
                            if not paragraph.runs:
                                continue

                            # Combine full text safely
                            full_text = ''
                            run_map = []  # (run, start_idx, end_idx)
                            for run in paragraph.runs:
                                text = run.text or ''
                                start_idx = len(full_text)
                                full_text += text
                                end_idx = len(full_text)
                                run_map.append((run, start_idx, end_idx))

                            if search_str in full_text:
                                # Replace in full text
                                full_text = full_text.replace(search_str, str(repl_str))

                                # Clear all runs
                                for run in paragraph.runs:
                                    run.text = ''

                                # Put replaced text into the first run
                                if run_map:
                                    first_run = run_map[0][0]
                                    first_run.text = full_text
                                    first_run.font.name = "Lato"
                                    first_run.font.size = Pt(10)

def fill_slides(prs, quarterized_results, bill_results):
    Q_START = '2024-Q1'
    replacemytext(prs, "<CRPY>", format_number(quarterized_results.loc['CREDITS', Q_START], 0))
    QUARTER_END = '2025-Q1'
    replacemytext(prs, "<CRCY>", format_number(quarterized_results.loc['CREDITS', QUARTER_END], 0))
    replacemytext(prs, "<CR%>", str(format_number(quarterized_results.loc['CREDITS', 'YoY%'], 0)) + "%")

    replacemytext(prs, "<CR/TBPY>", format_number(quarterized_results.loc['CR_PER_TB', Q_START], 2))
    replacemytext(prs, "<CR/TBCY>", format_number(quarterized_results.loc['CR_PER_TB', QUARTER_END], 2))
    replacemytext(prs, "<CR/TB%>", str(format_number(quarterized_results.loc['CR_PER_TB', 'YoY%'], 0)) + "%")

    replacemytext(prs, "<CR/HPY>", format_number(quarterized_results.loc['CR_PER_H', Q_START], 2))
    replacemytext(prs, "<CR/HCY>", format_number(quarterized_results.loc['CR_PER_H', QUARTER_END], 2))
    replacemytext(prs, "<CR/H%>", str(format_number(quarterized_results.loc['CR_PER_H', 'YoY%'], 0)) + "%")

    replacemytext(prs, "<TBPY>", format_number(quarterized_results.loc['TB_SCANNED', Q_START], 0))
    replacemytext(prs, "<TBCY>", format_number(quarterized_results.loc['TB_SCANNED', QUARTER_END], 0))
    replacemytext(prs, "<TB%>", str(format_number(quarterized_results.loc['TB_SCANNED', 'YoY%'], 0)) + "%")

    replacemytext(prs, "<HPY>", format_number(quarterized_results.loc['CPU_HOURS', Q_START], 0))
    replacemytext(prs, "<HCY>", format_number(quarterized_results.loc['CPU_HOURS', QUARTER_END], 0))
    replacemytext(prs, "<H%>", str(format_number(quarterized_results.loc['CPU_HOURS', 'YoY%'], 0)) + "%")

    replacemytext(prs, "<JUPY>", format_number(quarterized_results.loc['NB_USERS', Q_START], 0))
    replacemytext(prs, "<JUCY>", format_number(quarterized_results.loc['NB_USERS', QUARTER_END], 0))
    replacemytext(prs, "<JU%>", str(format_number(quarterized_results.loc['NB_USERS', 'YoY%'], 0)) + "%")

    replacemytext(prs, "<SCPY>", format_number(quarterized_results.loc['NB_SCHEMA', Q_START], 0))
    replacemytext(prs, "<SCCY>", format_number(quarterized_results.loc['NB_SCHEMA', QUARTER_END], 0))
    replacemytext(prs, "<SC%>", str(format_number(quarterized_results.loc['NB_SCHEMA', 'YoY%'], 0)) + "%")

    replacemytext(prs, "<WHPY>", format_number(quarterized_results.loc['NB_WH', Q_START], 0))
    replacemytext(prs, "<WHCY>", format_number(quarterized_results.loc['NB_WH', QUARTER_END], 0))
    replacemytext(prs, "<WH%>", str(format_number(quarterized_results.loc['NB_WH', 'YoY%'], 0)) + "%")

    replacemytext(prs, "<CR/QPY>", format_number(quarterized_results.loc['CR_PER_1000_JOBS', Q_START], 2))
    replacemytext(prs, "<CR/QCY>", format_number(quarterized_results.loc['CR_PER_1000_JOBS', QUARTER_END], 2))
    replacemytext(prs, "<CR/Q%>", str(format_number(quarterized_results.loc['CR_PER_1000_JOBS', 'YoY%'], 0)) + "%")

    replacemytext(prs, "<QPY>", format_number(quarterized_results.loc['XP_JOBS', Q_START] / pow(10, 6), 1))
    replacemytext(prs, "<QCY>", format_number(quarterized_results.loc['XP_JOBS', QUARTER_END] / pow(10, 6), 1))
    replacemytext(prs, "<Q%>", str(format_number(quarterized_results.loc['XP_JOBS', 'YoY%'], 0)) + "%")

    first_row = bill_results.iloc[0]
    last_row = bill_results.iloc[-1]

    growth_compute = ((last_row['COMPUTE'] - first_row['COMPUTE']) / first_row['COMPUTE']) * 100
    growth_storage = ((last_row['STORAGE'] - first_row['STORAGE']) / first_row['STORAGE']) * 100
    growth_other = ((last_row['OTHERS'] - first_row['OTHERS']) / first_row['OTHERS']) * 100

    replacemytext(prs, "<CG%>", str(format_number(growth_compute, 0)) + "%")
    replacemytext(prs, "<SG%>", str(format_number(growth_storage, 0)) + "%")
    replacemytext(prs, "<OG%>", str(format_number(growth_other, 0)) + "%")

    print(bill_results)
    l3m_df = bill_results.tail(3)
    sum_last_three_months = l3m_df[['COMPUTE', 'STORAGE', 'OTHERS']].sum()
    total_sum_last_three_months = sum_last_three_months.sum()
    percentage_contribution = (sum_last_three_months / total_sum_last_three_months) * 100

    # 5. Store each percentage in a separate variable
    compute_contribution = percentage_contribution['COMPUTE']
    storage_contribution = percentage_contribution['STORAGE']
    others_contribution = percentage_contribution['OTHERS']

    replacemytext(prs, "<CS%>", str(format_number(compute_contribution, 0)) + "%")
    replacemytext(prs, "<SS%>", str(format_number(storage_contribution, 0)) + "%")
    replacemytext(prs, "<OS%>", str(format_number(others_contribution, 0)) + "%")


    fill_chart(prs, quarterized_results, bill_results)
    pass


def fill_chart(prs, quarterized_results, bill_results):
    df_copy = quarterized_results.copy()
    df_yoy= quarterized_results.copy()

    df_copy = df_copy.drop(columns=['YoY%'])

    chart_data = CategoryChartData()
    chart_data.categories = df_copy.columns.values
    chart_data.add_series("# of TBs ", (df_copy.loc['TB_SCANNED']))
    chart_data.add_series("Cred / TB", (df_copy.loc['CR_PER_TB']))
    print(chart_data)

    # Delete the last column by name
    for shape in prs.slides[5].shapes:
        if shape.has_chart:
            chart = shape.chart
            print(chart.chart_type)
            chart.replace_data(chart_data)

    chart_data = CategoryChartData()
    chart_data.categories = df_copy.columns.values
    chart_data.add_series("Jobs (in M)", (df_copy.loc['XP_JOBS'] / pow(10, 6)))
    chart_data.add_series("Cred / 1000 J", (df_copy.loc['CR_PER_1000_JOBS']))
    print(chart_data)

    # Delete the last column by name
    for shape in prs.slides[3].shapes:
        if shape.has_chart:
            chart = shape.chart
            print(chart.chart_type)
            chart.replace_data(chart_data)

    ######################################################
    #################################################

    chart_data = CategoryChartData()
    chart_data.categories = df_copy.columns.values
    chart_data.add_series("# of Hrs ", (df_copy.loc['CPU_HOURS']))
    chart_data.add_series("Cred / Hr", (df_copy.loc['CR_PER_H']))
    try:
        # Delete the last column by name
        for shape in prs.slides[7].shapes:
            if shape.has_chart:
                chart = shape.chart
                print(chart.chart_type)
                chart.replace_data(chart_data)
    except:
        print("error")

    ######################################################
    #################################################

    columns_to_drop = ['2023-Q1', '2023-Q2', '2023-Q3', '2023-Q4']
    existing_columns_to_drop = [col for col in columns_to_drop if col in df_copy.columns]
    df_copy = df_copy.drop(columns=existing_columns_to_drop)

    chart_data = CategoryChartData()
    chart_data.categories = df_copy.columns.values
    chart_data.add_series("Jobs (in M) ", (df_copy.loc['XP_JOBS'] / pow(10, 6)))
    chart_data.add_series("Cred / 1000 J", (df_copy.loc['CR_PER_1000_JOBS']))
    # Delete the last column by name
    for shape in prs.slides[4].shapes:
        if shape.has_chart:
            chart = shape.chart
            print(chart.chart_type)
            chart.replace_data(chart_data)

    ######################################################
    #################################################
    chart_data = CategoryChartData()
    chart_data.categories = df_copy.columns.values
    chart_data.add_series("# of TBs ", (df_copy.loc['TB_SCANNED']))
    chart_data.add_series("Cred / TB", (df_copy.loc['CR_PER_TB']))
    # Delete the last column by name
    for shape in prs.slides[6].shapes:
        if shape.has_chart:
            chart = shape.chart
            chart.replace_data(chart_data)

    ######################################################
    #################################################
    chart_data = CategoryChartData()
    chart_data.categories = df_copy.columns.values
    chart_data.add_series("# of Hrs ", (df_copy.loc['CPU_HOURS']))
    chart_data.add_series("Cred / Hr", (df_copy.loc['CR_PER_H']))
    # Delete the last column by name
    for shape in prs.slides[8].shapes:
        if shape.has_chart:
            chart = shape.chart
            print(chart.chart_type)
            chart.replace_data(chart_data)

    ######################################################
    #################################################

    if 'COMPUTE' in bill_results.index:
        bill_results = bill_results.T

    # Create new chart data
    chart_data = CategoryChartData()
    chart_data.categories = bill_results.MONTH.tolist()  # dates as x-axis


    # Add each series
    chart_data.add_series('Compute', bill_results['COMPUTE'].tolist())
    chart_data.add_series('Storage', bill_results['STORAGE'].tolist())
    chart_data.add_series('Others', bill_results['OTHERS'].tolist())

    # Replace chart data on slide 10
    for shape in prs.slides[9].shapes:
        if shape.has_chart:
            chart = shape.chart
            print("Replacing chart of type:", chart.chart_type)
            chart.replace_data(chart_data)
            break  # remove this if you expect multiple charts to update


    data = {
        "Change between 2024-Q1 & 2025-Q1": [
            round(df_yoy.loc['CREDITS']['YoY%'],0),
            round(df_yoy.loc['XP_JOBS']['YoY%'],0),
            round(df_yoy.loc['TB_SCANNED']['YoY%'],0),
            round( df_yoy.loc['CPU_HOURS']['YoY%'],0),
            round(df_yoy.loc['NB_USERS']['YoY%'],0),
            round(df_yoy.loc['NB_SCHEMA']['YoY%'],0),
            round(df_yoy.loc['NB_WH']['YoY%'],0)
        ]
    }
    # Create DataFrame
    df_columns = pd.DataFrame(data, index=[
        "Compute Spend",
        "XP Jobs",
        "TB Scanned",
        "DURATION (H)",
        "# of Job Users",
        "# of Schemas",
        "# of Warehouses"
    ])

    df_columns["Change between 2024-Q1 & 2025-Q1"] = df_columns["Change between 2024-Q1 & 2025-Q1"]

    chart_data = CategoryChartData()
    chart_data.categories = df_columns.index.tolist()

    chart_data.add_series("Change between 2024-Q1 & 2025-Q1", df_columns['Change between 2024-Q1 & 2025-Q1'].tolist())
    try:
        # Delete the last column by name
        for shape in prs.slides[2].shapes:
            print(f"Shape : Type = {shape.shape_type}")

            if shape.has_chart:
                chart = shape.chart
                print(chart.chart_type )
                chart.replace_data(chart_data)
    except Exception as e:
        print(e)


def edit_ppt(ve_name, company_name, quarterized_results, bill_results):
    today = date.today()
    month = today.strftime("%B %Y")
    prs = Presentation('FINOPS_DRAFT.pptx')
    slides = prs.slides

    common_slides(company_name, month, prs, ve_name)

    fill_slides(prs, quarterized_results, bill_results)

    folder_name = "output/" + ve_name + "/" + company_name
    if not os.path.exists(folder_name):
        # If it does not exist, create it
        os.makedirs(folder_name)

    generated_ue_ppt = company_name + " - " + month + ".pptx"
    generated_ue_ppt = os.path.join(folder_name, generated_ue_ppt)
    # generated_ue_ppt = os.path.join(
    #      """/Users/siarora/Library/CloudStorage/GoogleDrive-simar.arora@snowflake.com/My Drive""",   generated_ue_ppt)
    # print(os.path.abspath(generated_ue_ppt))
    prs.save(generated_ue_ppt)
    return generated_ue_ppt


def calculateYoY(quarterly_metrics):
    Q_START = ('2024-Q1') ##todo
    QUARTER_END = '2025-Q1' ##todo quarter view
    quarterly_metrics['YoY%'] = ((quarterly_metrics[QUARTER_END] - quarterly_metrics[Q_START]) / quarterly_metrics[
        Q_START]) * 100

    return quarterly_metrics


def get_df_quarterly(wh_df, ue_df, ue_q_df):
    ue_df['CREDITS'] = 0

    for index, row in ue_df.iterrows():
        condition = (wh_df['MONTH'] == row['MONTH']) & (wh_df['ACCOUNT_ID'] == row['ACCOUNT_ID']) & (
                wh_df['ENTITY_ID'] == row['WAREHOUSE_ID']) & (wh_df['ENTITY_NAME'] == row['WAREHOUSE_NAME'])
        ue_df.at[index, 'CREDITS'] = wh_df.loc[condition, 'CREDITS_XP'].sum()

    quarterly_metrics = ue_df.groupby('PERIOD')[[
        'XP_JOBS', 'BYTES_SCANNED', 'CPU_HOURS', 'CREDITS']].sum()

    ue_q_quarterly_metrics = ue_q_df.groupby('PERIOD')[[
        'NB_DB', 'NB_SCHEMA', 'NB_USERS', 'NB_WH']].sum()

    quarterly_metrics['TB_SCANNED'] = quarterly_metrics['BYTES_SCANNED'] / pow(1024, 4)
    quarterly_metrics.drop('BYTES_SCANNED', axis=1, inplace=True)

    quarterly_metrics['NB_SCHEMA'] = ue_q_quarterly_metrics['NB_SCHEMA']
    # quarterly_metrics['NB_SCHEMA'] = quarterly_metrics['NB_SCHEMA'].astype(int)

    quarterly_metrics['NB_USERS'] = ue_q_quarterly_metrics['NB_USERS']
    # quarterly_metrics['NB_USERS'] = quarterly_metrics['NB_USERS'].astype(int)

    quarterly_metrics['NB_WH'] = ue_q_quarterly_metrics['NB_WH']
    # quarterly_metrics['NB_WH'] = quarterly_metrics['NB_WH'].astype(int)

    quarterly_metrics['CREDITS'] = quarterly_metrics['CREDITS'].apply(float)
    quarterly_metrics['TB_SCANNED'] = quarterly_metrics['TB_SCANNED'].apply(float)
    quarterly_metrics['CPU_HOURS'] = quarterly_metrics['CPU_HOURS'].apply(float)

    quarterly_metrics['CR_PER_1000_JOBS'] = quarterly_metrics['CREDITS'] * 1000 / quarterly_metrics['XP_JOBS']
    quarterly_metrics['CR_PER_TB'] = quarterly_metrics['CREDITS'] / quarterly_metrics['TB_SCANNED']
    quarterly_metrics['CR_PER_H'] = quarterly_metrics['CREDITS'] / quarterly_metrics['CPU_HOURS']
    quarterly_metrics['JOBS_PER_H'] = quarterly_metrics['XP_JOBS'] / quarterly_metrics['CPU_HOURS']

    transposed_quarterly_metrics = calculateYoY(quarterly_metrics.T)

    print(transposed_quarterly_metrics)

    return transposed_quarterly_metrics


def get_data(cursor_TAM, salesforce_id, salesforce_name, ve_name, cached=0):

    if cached:
        wh_level_credits_list = retrieve_my_list('wh_credits_list')
        unit_economics_wh_level_list = retrieve_my_list('ue_list')
        unit_economics_q_level_list = retrieve_my_list('ue_q_list')
    else:
        ue_q_level_queries = get_deployment_zone_queries_quarter_level(cursor_TAM, salesforce_id)
        cache_my_list(fetch_results_list_of_queries(cursor_TAM, ue_q_level_queries, 1), 'ue_q_list')
        unit_economics_q_level_list = retrieve_my_list('ue_q_list')

        cache_my_list(fetch_results(cursor_TAM, get_warehouse_level_query(salesforce_id)), 'wh_credits_list')
        wh_level_credits_list = retrieve_my_list('wh_credits_list')

        ue_wh_level_queries = get_deployment_zone_queries_wh_level(cursor_TAM, salesforce_id)
        cache_my_list(fetch_results_list_of_queries(cursor_TAM, ue_wh_level_queries, 1), 'ue_list')
        unit_economics_wh_level_list = retrieve_my_list('ue_list')
    filename = "output/" + ve_name + "/" + salesforce_name + "/" + "ue_q_list.csv"

    # Open the file in write mode
    with open(filename, mode='w', newline='') as file:
        writer = csv.writer(file)

        # Write the list of lists to the CSV file
        writer.writerows(unit_economics_q_level_list)

    filename = "output/" + ve_name + "/" + salesforce_name + "/" + "wh_level.csv"

    # Open the file in write mode
    with open(filename, mode='w', newline='') as file:
        writer = csv.writer(file)

        # Write the list of lists to the CSV file
        writer.writerows(wh_level_credits_list)

    filename = "output/" + ve_name + "/" + salesforce_name + "/" + "ue_list.csv"

    # Open the file in write mode
    with open(filename, mode='w', newline='') as file:
        writer = csv.writer(file)

        # Write the list of lists to the CSV file
        writer.writerows(unit_economics_wh_level_list)

    wh_df = pd.DataFrame(wh_level_credits_list[1:], columns=wh_level_credits_list[0])
    ue_df = pd.DataFrame(unit_economics_wh_level_list[1:], columns=unit_economics_wh_level_list[0])
    ue_q_df = pd.DataFrame(unit_economics_q_level_list[1:], columns=unit_economics_q_level_list[0])

    return wh_df, ue_df, ue_q_df


def get_quarterised_results(cursor_TAM, salesforce_id, salesforce_name, ve_name):
    wh_df, ue_df, ue_q_df = get_data(cursor_TAM, salesforce_id, salesforce_name, ve_name, cached=0)
    quarterly_metrics = get_df_quarterly(wh_df, ue_df, ue_q_df)
    #quarterly_metrics.to_csv(salesforce_name + '.csv', index=False)

    return quarterly_metrics


def get_bill(cursor_TAM, salesforce_id, salesforce_name):
    bill_query = BILLING.format(salesforce_id)
    print(bill_query)

    results = cursor_TAM.execute(bill_query).fetchall()
    headers = ["MONTH", "COMPUTE", "STORAGE", "OTHERS"]
    billing_df = pd.DataFrame(results[0:], columns=headers)

    return billing_df


def call_unit_economics(row):
    cursor_TAM = set_connection("SALES_MODELING_RL", "FINOPS_WH")

    ve_name, salesforce_name, salesforce_id = row
    print(ve_name, salesforce_name, salesforce_id)
    folder_name = "output/" + ve_name + "/" + salesforce_name
    if not os.path.exists(folder_name):
        # If it does not exist, create it
        os.makedirs(folder_name)
    quarterized_results = get_quarterised_results(cursor_TAM, salesforce_id, salesforce_name, ve_name)
    bill_results = get_bill(cursor_TAM, salesforce_id, salesforce_name)
    try:
        return edit_ppt(ve_name, salesforce_name, quarterized_results, bill_results)
    except Exception as e:
        print("Failed for " + salesforce_name + e)


if __name__ == '__main__':
    # cursor_SE = set_connection("SALES_ENGINEER", "SNOWADHOC")
    cursor_TAM = set_connection("SALES_MODELING_RL", "FINOPS_WH")

    with open('uniteconomics/customers.csv', mode='r') as file:
        csv_reader = csv.reader(file)
        for row in csv_reader:
            call_unit_economics(row)
